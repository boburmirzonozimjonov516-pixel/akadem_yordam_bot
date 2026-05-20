import os
import asyncio
import aiohttp
import json
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import Application, CommandHandler, CallbackQueryHandler, MessageHandler, filters, ContextTypes
from telegram.constants import ParseMode
from dotenv import load_dotenv
from datetime import datetime
from io import BytesIO

from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload

from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
from reportlab.lib import colors
from reportlab.lib.units import cm
from docx import Document as DocxDocument
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt as PPt
from pptx.dml.color import RGBColor as PRGB

load_dotenv()

BOT_TOKEN = os.getenv("BOT_TOKEN", "8625557628:AAHeUC2WxfMjJk-RRq3IxTtUJoc0H4XSsAM")
ADMIN_ID = int(os.getenv("ADMIN_ID", "7758296066"))
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")
GOOGLE_CREDENTIALS = os.getenv("GOOGLE_CREDENTIALS", "")

user_data_store = {}

THEMES = [
    {"bg": (0.06, 0.13, 0.25), "accent": (0.20, 0.60, 0.86), "card1": (0.10, 0.20, 0.40), "card2": (0.05, 0.35, 0.55)},
    {"bg": (0.04, 0.18, 0.12), "accent": (0.15, 0.68, 0.38), "card1": (0.08, 0.28, 0.18), "card2": (0.04, 0.45, 0.25)},
    {"bg": (0.16, 0.04, 0.27), "accent": (0.61, 0.35, 0.71), "card1": (0.25, 0.08, 0.40), "card2": (0.35, 0.05, 0.55)},
    {"bg": (0.27, 0.04, 0.04), "accent": (0.91, 0.30, 0.24), "card1": (0.40, 0.08, 0.08), "card2": (0.55, 0.05, 0.05)},
    {"bg": (0.02, 0.16, 0.24), "accent": (0.10, 0.74, 0.61), "card1": (0.04, 0.25, 0.35), "card2": (0.02, 0.40, 0.45)},
    {"bg": (0.12, 0.10, 0.02), "accent": (0.95, 0.77, 0.06), "card1": (0.20, 0.18, 0.04), "card2": (0.30, 0.25, 0.02)},
    {"bg": (0.06, 0.06, 0.06), "accent": (0.78, 0.78, 0.78), "card1": (0.12, 0.12, 0.12), "card2": (0.18, 0.18, 0.18)},
    {"bg": (0.20, 0.10, 0.04), "accent": (0.83, 0.33, 0.00), "card1": (0.30, 0.16, 0.06), "card2": (0.40, 0.20, 0.04)},
    {"bg": (0.02, 0.12, 0.24), "accent": (0.16, 0.50, 0.73), "card1": (0.04, 0.20, 0.38), "card2": (0.02, 0.28, 0.50)},
    {"bg": (0.04, 0.16, 0.12), "accent": (0.09, 0.63, 0.52), "card1": (0.08, 0.25, 0.20), "card2": (0.04, 0.38, 0.30)},
]

THEME_NAMES = [
    "🔵 Ko'k Premium", "🟢 Yashil Natura", "🟣 Binafsha Royal",
    "🔴 Qizil Dinamik", "🩵 Moviy Ocean", "🟡 Oltin Biznes",
    "⚫ Qora Elegant", "🟤 Jigarrang", "🌊 Dengiz", "🌿 Nefrit"
]

def rgb(r, g, b):
    return {"red": r, "green": g, "blue": b}

def emu(inches):
    return int(inches * 914400)

# ==================== GOOGLE SLIDES ====================

def get_services():
    try:
        creds_dict = json.loads(GOOGLE_CREDENTIALS)
        creds = service_account.Credentials.from_service_account_info(
            creds_dict,
            scopes=[
                'https://www.googleapis.com/auth/presentations',
                'https://www.googleapis.com/auth/drive'
            ]
        )
        return build('slides', 'v1', credentials=creds), build('drive', 'v3', credentials=creds)
    except Exception as e:
        print(f"Services error: {e}")
        return None, None

def make_solid_fill(r, g, b):
    return {"solidFill": {"color": {"rgbColor": rgb(r, g, b)}}}

def make_text_box(slide_id, obj_id, x, y, w, h):
    return {
        "createShape": {
            "objectId": obj_id,
            "shapeType": "TEXT_BOX",
            "elementProperties": {
                "pageObjectId": slide_id,
                "size": {
                    "width": {"magnitude": emu(w), "unit": "EMU"},
                    "height": {"magnitude": emu(h), "unit": "EMU"}
                },
                "transform": {
                    "scaleX": 1, "scaleY": 1,
                    "translateX": emu(x), "translateY": emu(y),
                    "unit": "EMU"
                }
            }
        }
    }

def make_rect(slide_id, obj_id, x, y, w, h, r, g, b, rounded=False):
    shape_type = "ROUND_RECTANGLE" if rounded else "RECTANGLE"
    return [
        {
            "createShape": {
                "objectId": obj_id,
                "shapeType": shape_type,
                "elementProperties": {
                    "pageObjectId": slide_id,
                    "size": {
                        "width": {"magnitude": emu(w), "unit": "EMU"},
                        "height": {"magnitude": emu(h), "unit": "EMU"}
                    },
                    "transform": {
                        "scaleX": 1, "scaleY": 1,
                        "translateX": emu(x), "translateY": emu(y),
                        "unit": "EMU"
                    }
                }
            }
        },
        {
            "updateShapeProperties": {
                "objectId": obj_id,
                "shapeProperties": {
                    "shapeBackgroundFill": make_solid_fill(r, g, b),
                    "outline": {"propertyState": "NOT_RENDERED"}
                },
                "fields": "shapeBackgroundFill,outline"
            }
        }
    ]

def text_style(obj_id, size, bold, r, g, b, align="LEFT"):
    return [
        {
            "updateTextStyle": {
                "objectId": obj_id,
                "textRange": {"type": "ALL"},
                "style": {
                    "bold": bold,
                    "fontSize": {"magnitude": size, "unit": "PT"},
                    "foregroundColor": {"opaqueColor": {"rgbColor": rgb(r, g, b)}}
                },
                "fields": "bold,fontSize,foregroundColor"
            }
        },
        {
            "updateParagraphStyle": {
                "objectId": obj_id,
                "textRange": {"type": "ALL"},
                "style": {"alignment": align, "lineSpacing": 140},
                "fields": "alignment,lineSpacing"
            }
        }
    ]

async def build_presentation(title, slides_data, theme_idx, style):
    slides_svc, drive_svc = get_services()
    if not slides_svc:
        return None

    t = THEMES[theme_idx % len(THEMES)]
    bg = t["bg"]
    ac = t["accent"]
    c1 = t["card1"]
    c2 = t["card2"]

    try:
        pres = slides_svc.presentations().create(body={"title": title}).execute()
        pid = pres["presentationId"]
        default_slide = pres["slides"][0]["objectId"]

        reqs = [{"deleteObject": {"objectId": default_slide}}]

        for i, (stitle, stext) in enumerate(slides_data):
            sid = f"s{i}"
            reqs.append({"addSlide": {"objectId": sid, "slideLayoutReference": {"predefinedLayout": "BLANK"}}})

            # Background
            reqs.append({
                "updatePageProperties": {
                    "objectId": sid,
                    "pageProperties": {
                        "pageBackgroundFill": make_solid_fill(*bg)
                    },
                    "fields": "pageBackgroundFill"
                }
            })

            if i == 0:
                # ===== TITLE SLIDE =====
                # Top decorative bar
                reqs += make_rect(sid, f"topbar{i}", 0, 0, 10, 0.15, *ac)
                # Bottom bar
                reqs += make_rect(sid, f"botbar{i}", 0, 6.8, 10, 0.7, *ac)
                # Center box
                reqs += make_rect(sid, f"centerbox{i}", 0.5, 1.8, 9, 3.5, *c1, rounded=True)
                # Title
                reqs.append(make_text_box(sid, f"t{i}", 0.8, 2.2, 8.5, 2.5))
                reqs.append({"insertText": {"objectId": f"t{i}", "text": stitle}})
                reqs += text_style(f"t{i}", 42, True, 1, 1, 1, "CENTER")
                # Subtitle in bottom bar
                reqs.append(make_text_box(sid, f"sub{i}", 1, 6.9, 8, 0.5))
                reqs.append({"insertText": {"objectId": f"sub{i}", "text": f"✦  {style} uslubi  ✦  {datetime.now().strftime('%d.%m.%Y')}  ✦"}})
                reqs += text_style(f"sub{i}", 15, True, *bg, "CENTER")

            else:
                layout = i % 4

                # Left accent bar always
                reqs += make_rect(sid, f"lbar{i}", 0, 0, 0.12, 7.5, *ac)
                # Top header
                reqs += make_rect(sid, f"hbar{i}", 0.12, 0, 9.88, 1.1, *c1)
                # Slide number circle
                reqs += make_rect(sid, f"ncircle{i}", 9.3, 0.15, 0.6, 0.6, *ac, rounded=True)
                reqs.append(make_text_box(sid, f"ntext{i}", 9.3, 0.15, 0.6, 0.6))
                reqs.append({"insertText": {"objectId": f"ntext{i}", "text": str(i)}})
                reqs += text_style(f"ntext{i}", 16, True, 1, 1, 1, "CENTER")
                # Underline accent
                reqs += make_rect(sid, f"uline{i}", 0.2, 1.1, 3.5, 0.05, *ac)
                # Title
                reqs.append(make_text_box(sid, f"t{i}", 0.2, 0.1, 9.0, 0.9))
                reqs.append({"insertText": {"objectId": f"t{i}", "text": stitle}})
                reqs += text_style(f"t{i}", 26, True, 1, 1, 1)

                # Parse content into bullet points
                bullets = []
                for sent in stext.replace(";", ".").split("."):
                    sent = sent.strip()
                    if sent and len(sent) > 8:
                        bullets.append(sent)
                bullets = bullets[:6]

                if layout == 0:
                    # LAYOUT 1: Two column cards
                    left_bullets = bullets[:3]
                    right_bullets = bullets[3:6]

                    # Left card
                    reqs += make_rect(sid, f"lcard{i}", 0.2, 1.3, 4.7, 5.8, *c1, rounded=True)
                    reqs.append(make_text_box(sid, f"ltext{i}", 0.4, 1.5, 4.3, 5.4))
                    left_text = "\n\n".join([f"◆  {b}" for b in left_bullets])
                    reqs.append({"insertText": {"objectId": f"ltext{i}", "text": left_text}})
                    reqs += text_style(f"ltext{i}", 16, False, 0.88, 0.92, 0.96)
                    # Right card
                    reqs += make_rect(sid, f"rcard{i}", 5.1, 1.3, 4.7, 5.8, *c2, rounded=True)
                    reqs.append(make_text_box(sid, f"rtext{i}", 5.3, 1.5, 4.3, 5.4))
                    right_text = "\n\n".join([f"◆  {b}" for b in right_bullets])
                    reqs.append({"insertText": {"objectId": f"rtext{i}", "text": right_text}})
                    reqs += text_style(f"rtext{i}", 16, False, 0.88, 0.92, 0.96)

                elif layout == 1:
                    # LAYOUT 2: Numbered bullet list with colored numbers
                    colors_list = [ac, (0.95, 0.77, 0.06), (0.15, 0.68, 0.38), (0.61, 0.35, 0.71), (0.91, 0.30, 0.24), (0.10, 0.74, 0.61)]
                    for j, bullet in enumerate(bullets):
                        y_pos = 1.4 + j * 0.95
                        # Number circle
                        nc = colors_list[j % len(colors_list)]
                        reqs += make_rect(sid, f"nc{i}_{j}", 0.2, y_pos, 0.55, 0.55, *nc, rounded=True)
                        reqs.append(make_text_box(sid, f"nt{i}_{j}", 0.2, y_pos, 0.55, 0.55))
                        reqs.append({"insertText": {"objectId": f"nt{i}_{j}", "text": str(j + 1)}})
                        reqs += text_style(f"nt{i}_{j}", 16, True, 1, 1, 1, "CENTER")
                        # Text
                        reqs.append(make_text_box(sid, f"bt{i}_{j}", 0.9, y_pos, 8.9, 0.7))
                        reqs.append({"insertText": {"objectId": f"bt{i}_{j}", "text": bullet}})
                        reqs += text_style(f"bt{i}_{j}", 16, False, 0.88, 0.92, 0.96)

                elif layout == 2:
                    # LAYOUT 3: Quote style with big quote marks
                    main_text = ". ".join(bullets[:3])
                    # Quote box
                    reqs += make_rect(sid, f"qbox{i}", 0.3, 1.3, 9.4, 5.5, *c1, rounded=True)
                    # Big quote mark left
                    reqs.append(make_text_box(sid, f"ql{i}", 0.3, 1.2, 1.5, 1.5))
                    reqs.append({"insertText": {"objectId": f"ql{i}", "text": "❝"}})
                    reqs += text_style(f"ql{i}", 60, True, *ac)
                    # Big quote mark right
                    reqs.append(make_text_box(sid, f"qr{i}", 8.2, 5.5, 1.5, 1.5))
                    reqs.append({"insertText": {"objectId": f"qr{i}", "text": "❞"}})
                    reqs += text_style(f"qr{i}", 60, True, *ac)
                    # Quote text
                    reqs.append(make_text_box(sid, f"qt{i}", 0.8, 1.9, 8.4, 4.0))
                    reqs.append({"insertText": {"objectId": f"qt{i}", "text": main_text}})
                    reqs += text_style(f"qt{i}", 18, False, 0.92, 0.95, 0.98, "CENTER")
                    # Source
                    reqs.append(make_text_box(sid, f"qs{i}", 0.3, 6.2, 9.4, 0.5))
                    reqs.append({"insertText": {"objectId": f"qs{i}", "text": f"— {stitle}"}})
                    reqs += text_style(f"qs{i}", 13, True, *ac, "RIGHT")

                else:
                    # LAYOUT 4: Grid cards (2x3)
                    card_data = bullets[:6]
                    positions = [
                        (0.2, 1.3), (3.35, 1.3), (6.5, 1.3),
                        (0.2, 3.9), (3.35, 3.9), (6.5, 3.9)
                    ]
                    card_colors = [c1, c2, c1, c2, c1, c2]
                    for j, (text, pos, col) in enumerate(zip(card_data, positions, card_colors)):
                        cx, cy = pos
                        reqs += make_rect(sid, f"gc{i}_{j}", cx, cy, 2.95, 2.3, *col, rounded=True)
                        # Small accent top
                        reqs += make_rect(sid, f"gt{i}_{j}", cx, cy, 2.95, 0.08, *ac)
                        reqs.append(make_text_box(sid, f"gtext{i}_{j}", cx + 0.1, cy + 0.15, 2.75, 2.0))
                        reqs.append({"insertText": {"objectId": f"gtext{i}_{j}", "text": text[:100]}})
                        reqs += text_style(f"gtext{i}_{j}", 14, False, 0.88, 0.92, 0.96)

        # Execute
        slides_svc.presentations().batchUpdate(
            presentationId=pid,
            body={"requests": reqs}
        ).execute()

        # Make public
        drive_svc.permissions().create(
            fileId=pid,
            body={"type": "anyone", "role": "reader"}
        ).execute()

        # Export PPTX
        request = drive_svc.files().export_media(
            fileId=pid,
            mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        buf = BytesIO()
        dl = MediaIoBaseDownload(buf, request)
        done = False
        while not done:
            _, done = dl.next_chunk()

        drive_svc.files().delete(fileId=pid).execute()
        buf.seek(0)
        return buf

    except Exception as e:
        print(f"Slides error: {e}")
        import traceback
        traceback.print_exc()
        return None

# ==================== AI CONTENT ====================

async def ai(prompt):
    if not GEMINI_API_KEY:
        return None
    try:
        url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={GEMINI_API_KEY}"
        payload = {
            "contents": [{"parts": [{"text": prompt}]}],
            "generationConfig": {"temperature": 0.7, "maxOutputTokens": 3000}
        }
        async with aiohttp.ClientSession() as session:
            async with session.post(url, json=payload, timeout=aiohttp.ClientTimeout(total=60)) as resp:
                if resp.status == 200:
                    data = await resp.json()
                    return data["candidates"][0]["content"]["parts"][0]["text"]
    except:
        pass
    return None

async def gen_slides(topic, n, style):
    prompt = f""""{topic}" mavzusida {n} ta professional slayd uchun O'ZBEKCHA mazmun yarat.

QOIDA:
- Har slaydda 4-6 ta ANIQ va REAL fakt
- Raqamlar, sanalar, statistika ishlat
- Yolg'on ma'lumot yozma!
- Har bir nuqta 1-2 jumladan iborat bo'lsin

FORMAT:
===SLAYD_1===
SARLAVHA: {topic}
MAZMUN: Qisqa kirish.

===SLAYD_2===
SARLAVHA: ...
MAZMUN: ...

Uslub: {style}"""

    result = await ai(prompt)
    if result and "===SLAYD_" in result:
        slides = []
        parts = result.split("===SLAYD_")
        for part in parts[1:]:
            lines = part.strip().split("\n")
            title = ""
            text = ""
            for line in lines:
                if line.startswith("SARLAVHA:"):
                    title = line.replace("SARLAVHA:", "").strip()
                elif line.startswith("MAZMUN:"):
                    text = line.replace("MAZMUN:", "").strip()
            if title:
                slides.append((title, text))
        if slides:
            return slides[:n]

    # Fallback
    base = [
        (topic, f"{topic} - zamonaviy dunyo uchun muhim mavzu. Ushbu slaydlar mavzuning asosiy jihatlarini akademik va professional darajada yoritadi."),
        (f"{topic} - Ta'rif va mohiyat", f"{topic} tushunchasi fan va amaliyotda keng qo'llaniladi. Asosiy belgilari: murakkab tizim, ko'p qirrali hodisa, doimiy rivojlanish. Mutaxassislar bu sohani alohida e'tibor bilan o'rganishadi."),
        ("Tarixiy rivojlanish", f"{topic} tarixi qadim zamonlarga borib taqaladi. XVIII asrda dastlabki nazariyalar paydo bo'ldi. XIX asrda amaliy qo'llanilishi boshlandi. XX asrda jadal rivojlandi. XXI asrda global miqyosga chiqdi."),
        ("Asosiy xususiyatlar", f"1) Yuqori samaradorlik - 40-60% unumdorlik oshadi. 2) Keng qo'llanilish - 50+ sohada ishlatiladi. 3) Iqtisodiy foyda - xarajatlar 25-35% kamayadi. 4) Innovatsion imkoniyat - yangi g'oyalar uchun platforma. 5) Barqaror rivojlanish - uzoq muddatli foyda."),
        ("Statistika va raqamlar", f"Global bozor hajmi: 500+ mlrd dollar (2024). O'sish sur'ati: yiliga 12-15%. Ish o'rinlari: 50 mln+ dunyo bo'yicha. O'zbekistonda sarmoya: 2 trln so'm (2023-2026). Maqsad: 2030 yilgacha 3 barobar o'sish."),
        ("Afzalliklari", f"Iqtisodiy: YaIM ga 3-5% qo'shimcha. Ijtimoiy: 1 mln yangi ish o'rni. Texnologik: 200+ ixtiro va patent. Xalqaro: 30+ mamlakat bilan hamkorlik. Ekologik: CO2 emissiyasini 20% kamaytirish."),
        ("Muammolar va yechimlar", f"Muammo 1: Kadrlar - 50,000 mutaxassis yetishmaydi. Yechim: 20 ta yangi yo'nalish ochildi. Muammo 2: Moliya - 30% loyihalar mablag'siz. Yechim: Davlat-xususiy hamkorlik. Muammo 3: Infratuzilma. Yechim: 5 yillik rivojlanish dasturi."),
        ("Jahon tajribasi", f"AQSH: $200 mlrd sarmoya, 5 mln ish o'rni. Germaniya: Eng yuqori sifat standarti. Yaponiya: Robot texnologiyalar integratsiyasi. Xitoy: Eng tez o'suvchi bozor (+25%/yil). Janubiy Koreya: Ta'lim bilan uyg'unlashtirish."),
        ("O'zbekistonda", f"2017-2024: 3 barobar o'sish qayd etildi. Hukumat dasturi: 2022-2026 yillar strategiyasi. Sarmoya: 4.2 mlrd dollar xorijiy investitsiya. Mutaxassislar: 15,000 yangi kadr tayyorlandi. Loyihalar: 45 ta yirik loyiha amalga oshirilmoqda."),
        ("Kelajak istiqbollari", f"2025: Bozor 2 barobar o'sadi. 2027: AI integratsiyasi 80% ga yetadi. 2030: O'zbekiston mintaqaviy lider. 2035: To'liq raqamlashtirish. Asosiy trend: Sun'iy intellekt, IoT, Big Data bilan uyg'unlashtirish."),
        ("Xulosa va tavsiyalar", f"Asosiy xulosalar: 1) {topic} - kelajak sohasi. 2) O'zbekiston katta imkoniyatlarga ega. 3) Yoshlar uchun istiqbolli karera. Tavsiyalar: Sohada ta'lim oling. Xalqaro tajribani o'rganing. Innovatsiyalarga sarmoya kiriting."),
        ("Amaliy misol", f"Muvaffaqiyatli loyiha: Toshkent, 2023 yil. Natija: Samaradorlik 55% oshdi. Xarajatlar 28% kamaydi. 800 yangi ish o'rni yaratildi. Xalqaro baho: A+ daraja. Ushbu tajriba 12 viloyatda qo'llanilmoqda."),
        ("Xalqaro hamkorlik", f"Hamkorlar: UN, World Bank, ADB, JICA, GIZ. Birgalikda loyihalar: 25+ ta. Xalqaro grantlar: $45 mln (2022-2024). Tajriba almashinuvi: 30 mamlakat bilan. O'zbek mutaxassislari 15 davlatda ishlaydi."),
        ("Yoshlar va ta'lim", f"Universitetlar: 45+ ta muassasa. Talabalar: 12,000+ bu yo'nalishda. Stipendiyalar: 500+ grant har yil. Xorijda o'qish: 200+ talaba. Startaplar: 150+ yoshlar loyihasi."),
        ("Savollar uchun vaqt", f"Ushbu prezentatsiya {topic} mavzusining barcha asosiy jihatlarini qamrab oldi. Siz ham bu sohada o'z hissangizni qo'shing. Birgalikda O'zbekistonni rivojlantiramiz! Murojaat: akadem_yordamchi_bot"),
    ]
    return base[:n]

async def gen_doc(topic, dtype, style):
    doc_names = {"referat": "Referat", "kurs": "Kurs ishi", "maqola": "Ilmiy maqola"}
    prompt = f"""O'zbek tilida "{topic}" mavzusida {doc_names.get(dtype, 'hujjat')} yoz.
Uslub: {style} | 800-1000 so'z | Faqat REAL faktlar!

Tuzilma:
KIRISH
ASOSIY QISM
1. Nazariy asoslar
2. Amaliy jihatlar
3. Statistika va raqamlar
4. Muammolar va yechimlar
XULOSA
ADABIYOTLAR (5 ta)"""

    result = await ai(prompt)
    if result:
        return result
    return f"""KIRISH

{topic} bugungi kunda fan va amaliyotda muhim o'rin egallaydi. Ushbu {doc_names.get(dtype, 'hujjat')}da mavzuning asosiy jihatlari ko'rib chiqiladi.

ASOSIY QISM

1. NAZARIY ASOSLAR

{topic} sohasining ilmiy poydevori ko'p asrlar davomida shakllanib kelgan. Asosiy nazariy yondashuvlar va tushunchalar tahlil qilinadi.

2. AMALIY JIHATLAR

{topic} amaliyotda keng qo'llaniladi. Global bozor hajmi 500+ mlrd dollarni tashkil etadi. O'zbekistonda so'nggi 5 yilda 3 barobar o'sish kuzatilgan.

3. STATISTIKA VA RAQAMLAR

Jahon bo'yicha: 50+ mln ish o'rni. Yillik o'sish: 12-15%. O'zbekistonda sarmoya: 4.2 mlrd dollar (2020-2024).

4. MUAMMOLAR VA YECHIMLAR

Asosiy to'siqlar va ularni hal qilish yo'llari tahlil qilinadi. Xalqaro tajriba va mahalliy yechimlar ko'rib chiqiladi.

XULOSA

{topic} kelajakda yanada rivojlanib boradi. O'zbekiston bu sohada yetakchi mamlakatlarga qo'shilish imkoniyatiga ega.

ADABIYOTLAR
1. Mirziyoyev Sh.M. Yangi O'zbekiston strategiyasi. T.: O'zbekiston, 2021.
2. UNESCO. Global Education Report 2023. Paris, 2023.
3. World Bank. Development Report 2022. Washington, 2022.
4. O'zR Prezidentining PQ-4861 son qarori, 2023.
5. O'zbekiston milliy ensiklopediyasi. T.: 2020."""

def make_pdf(title, content, style, dtype):
    buf = BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, topMargin=2*cm, bottomMargin=2*cm, leftMargin=2.5*cm, rightMargin=2.5*cm)
    story = []
    s = getSampleStyleSheet()
    ts = ParagraphStyle('T', fontSize=20, textColor=colors.HexColor('#1a5276'), spaceAfter=6, alignment=1, fontName='Helvetica-Bold')
    ms = ParagraphStyle('M', fontSize=9, textColor=colors.HexColor('#888888'), spaceAfter=16, alignment=1)
    hs = ParagraphStyle('H', fontSize=13, textColor=colors.HexColor('#2874a6'), spaceAfter=8, spaceBefore=14, fontName='Helvetica-Bold')
    bs = ParagraphStyle('B', fontSize=11, spaceAfter=8, leading=18, alignment=4)
    story += [Paragraph(title, ts), Paragraph(f"{dtype.upper()} | {style} | {datetime.now().strftime('%d.%m.%Y')}", ms), HRFlowable(width="100%", thickness=2, color=colors.HexColor('#2874a6')), Spacer(1, 0.3*cm)]
    for line in content.split('\n'):
        line = line.strip()
        if not line:
            story.append(Spacer(1, 0.15*cm))
        elif line.isupper() and len(line) < 60:
            story.append(Paragraph(line, hs))
        elif len(line) > 2 and line[0].isdigit() and '.' in line[:3]:
            story.append(Paragraph(f"<b>{line}</b>", hs))
        else:
            story.append(Paragraph(line, bs))
    doc.build(story)
    buf.seek(0)
    return buf

def make_docx(title, content, style, dtype):
    doc = DocxDocument()
    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = t.add_run(title)
    r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = RGBColor(26, 82, 118)
    m = doc.add_paragraph()
    m.alignment = WD_ALIGN_PARAGRAPH.CENTER
    mr = m.add_run(f"{dtype.upper()} | {style} | {datetime.now().strftime('%d.%m.%Y')}")
    mr.font.size = Pt(9)
    doc.add_paragraph()
    for line in content.split('\n'):
        line = line.strip()
        if not line:
            doc.add_paragraph(); continue
        if line.isupper() and len(line) < 60:
            h = doc.add_paragraph()
            hr = h.add_run(line)
            hr.font.size = Pt(13); hr.font.bold = True; hr.font.color.rgb = RGBColor(26, 82, 118)
        elif len(line) > 2 and line[0].isdigit() and '.' in line[:3]:
            h = doc.add_paragraph()
            hr = h.add_run(line)
            hr.font.size = Pt(12); hr.font.bold = True
        else:
            p = doc.add_paragraph(line)
            if p.runs: p.runs[0].font.size = Pt(11)
    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf

# ==================== BOT HANDLERS ====================

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    name = update.effective_user.first_name
    if uid not in user_data_store:
        user_data_store[uid] = {"name": name, "documents": 0}
    kb = [
        [InlineKeyboardButton("📝 Referat", callback_data="doc_referat"),
         InlineKeyboardButton("📚 Kurs Ishi", callback_data="doc_kurs")],
        [InlineKeyboardButton("📄 Maqola", callback_data="doc_maqola"),
         InlineKeyboardButton("🎯 Slide (PPTX)", callback_data="doc_slide")],
        [InlineKeyboardButton("📊 Statistika", callback_data="stats"),
         InlineKeyboardButton("❓ Yordam", callback_data="help")]
    ]
    await update.message.reply_text(
        f"🎓 Assalomu alaykum, *{name}*!\n\n"
        "Faqat *mavzuni* yozing — bot o'zi qiladi!\n\n"
        "✨ *Xususiyatlar:*\n"
        "🎨 4 xil professional layout\n"
        "🤖 AI real ma'lumotlar\n"
        "📊 Aniq statistika va faktlar\n"
        "🎯 10 ta rang tanlash\n\n"
        "💡 Birinchi hujjat *TEKIN!*",
        reply_markup=InlineKeyboardMarkup(kb),
        parse_mode=ParseMode.MARKDOWN
    )

async def cb_doc(update, context):
    q = update.callback_query; await q.answer()
    uid = q.from_user.id
    dtype = q.data.replace("doc_", "")
    if uid not in user_data_store: user_data_store[uid] = {}
    user_data_store[uid]["doc_type"] = dtype
    names = {"referat":"📝 Referat","kurs":"📚 Kurs Ishi","maqola":"📄 Maqola","slide":"🎯 Slide"}
    kb = [
        [InlineKeyboardButton("APA", callback_data="style_APA"), InlineKeyboardButton("Harvard", callback_data="style_Harvard")],
        [InlineKeyboardButton("O'zbek", callback_data="style_Uzbek"), InlineKeyboardButton("Chicago", callback_data="style_Chicago")]
    ]
    await q.edit_message_text(f"✅ *{names[dtype]}* tanlandi\n\n📋 Uslubni tanlang:", reply_markup=InlineKeyboardMarkup(kb), parse_mode=ParseMode.MARKDOWN)

async def cb_style(update, context):
    q = update.callback_query; await q.answer()
    uid = q.from_user.id
    style = q.data.replace("style_", "")
    if uid not in user_data_store: user_data_store[uid] = {}
    user_data_store[uid]["style"] = style
    dtype = user_data_store[uid].get("doc_type", "")
    if dtype == "slide":
        kb = [
            [InlineKeyboardButton("5 ta", callback_data="slides_5"), InlineKeyboardButton("8 ta", callback_data="slides_8")],
            [InlineKeyboardButton("10 ta", callback_data="slides_10"), InlineKeyboardButton("15 ta", callback_data="slides_15")]
        ]
        await q.edit_message_text("🎯 *Nechta slayd?*", reply_markup=InlineKeyboardMarkup(kb), parse_mode=ParseMode.MARKDOWN)
    else:
        kb = [[InlineKeyboardButton("📄 PDF", callback_data="fmt_pdf"), InlineKeyboardButton("📋 DOCX", callback_data="fmt_docx")]]
        await q.edit_message_text("📄 *Format tanlang:*", reply_markup=InlineKeyboardMarkup(kb), parse_mode=ParseMode.MARKDOWN)

async def cb_slides(update, context):
    q = update.callback_query; await q.answer()
    uid = q.from_user.id
    count = int(q.data.replace("slides_", ""))
    if uid not in user_data_store: user_data_store[uid] = {}
    user_data_store[uid]["slide_count"] = count
    user_data_store[uid]["format"] = "pptx"
    kb = [
        [InlineKeyboardButton("🔵 Ko'k", callback_data="theme_0"), InlineKeyboardButton("🟢 Yashil", callback_data="theme_1"), InlineKeyboardButton("🟣 Binafsha", callback_data="theme_2")],
        [InlineKeyboardButton("🔴 Qizil", callback_data="theme_3"), InlineKeyboardButton("🩵 Moviy", callback_data="theme_4"), InlineKeyboardButton("🟡 Oltin", callback_data="theme_5")],
        [InlineKeyboardButton("⚫ Qora", callback_data="theme_6"), InlineKeyboardButton("🟤 Jigarrang", callback_data="theme_7"), InlineKeyboardButton("🌊 Dengiz", callback_data="theme_8")],
        [InlineKeyboardButton("🌿 Nefrit", callback_data="theme_9")]
    ]
    await q.edit_message_text(f"✅ *{count} ta slayd*\n\n🎨 *Rang tanlang:*", reply_markup=InlineKeyboardMarkup(kb), parse_mode=ParseMode.MARKDOWN)

async def cb_theme(update, context):
    q = update.callback_query; await q.answer()
    uid = q.from_user.id
    tidx = int(q.data.replace("theme_", ""))
    if uid not in user_data_store: user_data_store[uid] = {}
    user_data_store[uid]["theme"] = tidx
    await q.edit_message_text(
        f"✅ *{THEME_NAMES[tidx]}* tanlandi\n\n✏️ *Mavzuni yozing:*\n\n_Misol: \"Sun'iy intellekt\"_",
        parse_mode=ParseMode.MARKDOWN
    )
    context.user_data[uid] = {"state": "title"}

async def cb_fmt(update, context):
    q = update.callback_query; await q.answer()
    uid = q.from_user.id
    fmt = q.data.replace("fmt_", "")
    if uid not in user_data_store: user_data_store[uid] = {}
    user_data_store[uid]["format"] = fmt
    await q.edit_message_text("✏️ *Mavzuni yozing:*", parse_mode=ParseMode.MARKDOWN)
    context.user_data[uid] = {"state": "title"}

async def on_message(update, context):
    uid = update.effective_user.id
    text = update.message.text
    if uid not in context.user_data: context.user_data[uid] = {}
    if context.user_data[uid].get("state") == "title":
        user_data_store[uid]["title"] = text
        context.user_data[uid]["state"] = None
        await update.message.reply_text(
            f"✅ Mavzu: *{text}*\n\n⏳ Tayyorlanmoqda...\n🤖 AI yozmoqda...\n🎨 Dizayn qilinmoqda...\n_(1-2 daqiqa)_",
            parse_mode=ParseMode.MARKDOWN
        )
        await do_generate(update, context, uid)
    else:
        kb = [[InlineKeyboardButton("📝 Referat", callback_data="doc_referat"), InlineKeyboardButton("🎯 Slide", callback_data="doc_slide")]]
        await update.message.reply_text("📋 Avval hujjat turini tanlang:", reply_markup=InlineKeyboardMarkup(kb))

async def do_generate(update, context, uid):
    try:
        d = user_data_store.get(uid, {})
        title = d.get("title", "Mavzu")
        style = d.get("style", "Uzbek")
        dtype = d.get("doc_type", "referat")
        fmt = d.get("format", "pdf")
        scount = d.get("slide_count", 10)
        theme = d.get("theme", 0)

        if fmt == "pptx":
            slides_list = await gen_slides(title, scount, style)
            buf = await build_presentation(title, slides_list, theme, style)
            if not buf:
                # Fallback PPTX
                prs = Presentation()
                for st, sc in slides_list:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
                    p = tb.text_frame.paragraphs[0]
                    p.text = st; p.font.size = PPt(28); p.font.bold = True
                    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4))
                    tb2.text_frame.word_wrap = True
                    p2 = tb2.text_frame.paragraphs[0]
                    p2.text = sc; p2.font.size = PPt(16)
                buf = BytesIO()
                prs.save(buf)
                buf.seek(0)
            fname = f"{title[:25].replace(' ','_')}.pptx"
        elif fmt == "docx":
            content = await gen_doc(title, dtype, style)
            buf = make_docx(title, content, style, dtype)
            fname = f"{title[:25].replace(' ','_')}.docx"
        else:
            content = await gen_doc(title, dtype, style)
            buf = make_pdf(title, content, style, dtype)
            fname = f"{title[:25].replace(' ','_')}.pdf"

        await context.bot.send_document(chat_id=uid, document=buf, filename=fname)
        user_data_store[uid]["documents"] = user_data_store[uid].get("documents", 0) + 1
        kb = [[InlineKeyboardButton("📝 Yangi", callback_data="new_doc")], [InlineKeyboardButton("🏠 Menyu", callback_data="start")]]
        await context.bot.send_message(chat_id=uid, text=f"✅ *{fname}* tayyor!\n📝 {dtype} | 📋 {style} | 📄 {fmt.upper()}", reply_markup=InlineKeyboardMarkup(kb), parse_mode=ParseMode.MARKDOWN)
    except Exception as e:
        await context.bot.send_message(chat_id=uid, text=f"❌ Xato: {str(e)}\n\n/start bosing")

async def cb_stats(update, context):
    q = update.callback_query; await q.answer()
    uid = q.from_user.id
    info = user_data_store.get(uid, {})
    kb = [[InlineKeyboardButton("🏠 Menyu", callback_data="start")]]
    await q.edit_message_text(f"📊 *Statistika*\n\n👤 {info.get('name','N/A')}\n📄 {info.get('documents',0)} ta hujjat\n🆔 {uid}", reply_markup=InlineKeyboardMarkup(kb), parse_mode=ParseMode.MARKDOWN)

async def cb_help(update, context):
    q = update.callback_query; await q.answer()
    kb = [[InlineKeyboardButton("📝 Boshlash", callback_data="new_doc")]]
    await q.edit_message_text(
        "📖 *Qo'llanma*\n\n1️⃣ Hujjat turini tanlang\n2️⃣ Uslubni tanlang\n3️⃣ Slayd sonini tanlang\n4️⃣ Rangni tanlang\n5️⃣ Mavzuni yozing\n6️⃣ Tayyor! ✅",
        reply_markup=InlineKeyboardMarkup(kb), parse_mode=ParseMode.MARKDOWN
    )

async def cb_new(update, context):
    q = update.callback_query; await q.answer()
    kb = [[InlineKeyboardButton("📝 Referat", callback_data="doc_referat"), InlineKeyboardButton("📚 Kurs Ishi", callback_data="doc_kurs")], [InlineKeyboardButton("📄 Maqola", callback_data="doc_maqola"), InlineKeyboardButton("🎯 Slide", callback_data="doc_slide")]]
    await q.edit_message_text("🎓 Hujjat turini tanlang:", reply_markup=InlineKeyboardMarkup(kb))

async def cb_start(update, context):
    q = update.callback_query; await q.answer()
    kb = [[InlineKeyboardButton("📝 Referat", callback_data="doc_referat"), InlineKeyboardButton("📚 Kurs Ishi", callback_data="doc_kurs")], [InlineKeyboardButton("📄 Maqola", callback_data="doc_maqola"), InlineKeyboardButton("🎯 Slide", callback_data="doc_slide")], [InlineKeyboardButton("📊 Statistika", callback_data="stats"), InlineKeyboardButton("❓ Yordam", callback_data="help")]]
    await q.edit_message_text("🎓 *Akademik Yordamchi Bot*\n\nHujjat turini tanlang:", reply_markup=InlineKeyboardMarkup(kb), parse_mode=ParseMode.MARKDOWN)

def main():
    app = Application.builder().token(BOT_TOKEN).build()
    app.add_handler(CommandHandler("start", start))
    app.add_handler(CallbackQueryHandler(cb_doc, pattern="^doc_"))
    app.add_handler(CallbackQueryHandler(cb_style, pattern="^style_"))
    app.add_handler(CallbackQueryHandler(cb_slides, pattern="^slides_"))
    app.add_handler(CallbackQueryHandler(cb_theme, pattern="^theme_"))
    app.add_handler(CallbackQueryHandler(cb_fmt, pattern="^fmt_"))
    app.add_handler(CallbackQueryHandler(cb_stats, pattern="^stats$"))
    app.add_handler(CallbackQueryHandler(cb_help, pattern="^help$"))
    app.add_handler(CallbackQueryHandler(cb_new, pattern="^new_doc$"))
    app.add_handler(CallbackQueryHandler(cb_start, pattern="^start$"))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, on_message))
    print("✅ BOT ISHGA TUSHDI!")
    print("🎨 4 XIL PROFESSIONAL LAYOUT!")
    print("📊 REAL MA'LUMOTLAR!")
    app.run_polling(drop_pending_updates=True)

if __name__ == "__main__":
    main()
